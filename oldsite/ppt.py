import os
import logging
from pptx import Presentation
from pptx.util import Inches



class Config:
    def __init__(self):
        self.slide_n = 0
        self.coord_list = []
        self.logger = 0
        self.valid = True
        self.check_files()
        self.get_info()
    def check_files(self):

        if not os.path.exists('logs'):
            os.makedirs('logs')

        logging.basicConfig(filename="logs\\outputlog.log",
                            format='%(asctime)s %(message)s',
                            filemode='w')
        self.logger = logging.getLogger()
        self.logger.setLevel(logging.DEBUG)




        if not os.path.exists('input'):
            os.makedirs('input')
            self.logger.error("The 'input' directory does not exist and has now been created.")
            self.valid = False

        if not os.path.exists('output'):
            os.makedirs('output')
            self.logger.error("The 'output' directory does not exist and has now been created.")
            self.valid = False

        if not os.path.exists('config'):
            os.makedirs('config')
            self.logger.error("The 'config' directory does not exist and has now been created.")
            self.valid = False




        # Using readlines()
        if not os.path.isfile('config\\PPT-XL.config'):
            self.logger.error("The 'PPT-XL.config' file does not exist and has now been created.")
            file1 = open('config\\PPT-XL.config', 'a')
            file1.close()
            self.valid = False
    def get_info(self):
        if self.valid:
            file1 = open('config\\PPT-XL.config', 'r')
            Lines = file1.readlines()
            self.slide_n = int(Lines[0])
             
            for c in range(len(Lines)):
                if Lines[c][0] == '#':
                    continue
                if c == 0:
                    continue
                print("Line{}: {}".format(c, Lines[c].strip()))

                self.coord_list.append(Lines[c].strip())
            file1.close()
        else:
            exit()

class Creator:
    def __init__(self):
        self.c = Config()
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

    def create_slide (self):

        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Hello, World!"
        self.prs.save('output\\test.pptx')

if __name__ == "__main__":
    c = Creator()
    c.create_ppt()
